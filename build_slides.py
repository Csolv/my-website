"""
Build ADC AI Project slides as a .pptx (Google Slides-importable)
Run: python3 build_slides.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colors ──────────────────────────────────────────────────────────────────
GREEN       = RGBColor(0x1B, 0x6B, 0x3A)   # ADC dark green
GREEN_LIGHT = RGBColor(0x2E, 0x9E, 0x58)   # accent green
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK        = RGBColor(0x12, 0x12, 0x12)
GRAY        = RGBColor(0x55, 0x55, 0x55)
LIGHT_BG    = RGBColor(0xF4, 0xF7, 0xF4)   # very light green-white

# ── Paths ────────────────────────────────────────────────────────────────────
BASE = "/Users/connorsolvason/my-website"
IMGS = {
    "hero":        f"{BASE}/Screenshot 2026-03-25 at 6.42.46 AM.png",
    "about":       f"{BASE}/Screenshot 2026-03-25 at 6.43.47 AM.png",
    "stats":       f"{BASE}/Screenshot 2026-03-25 at 6.44.53 AM.png",
    "bios":        f"{BASE}/Screenshot 2026-03-25 at 6.45.26 AM.png",
    "values":      f"{BASE}/Screenshot 2026-03-25 at 6.45.42 AM.png",
    "achievements":f"{BASE}/Screenshot 2026-03-25 at 6.45.49 AM.png",
    "calendar":    f"{BASE}/Screenshot 2026-03-25 at 6.46.32 AM.png",
    "participate": f"{BASE}/Screenshot 2026-03-25 at 6.47.01 AM.png",
    "contact":     f"{BASE}/Screenshot 2026-03-25 at 6.47.32 AM.png",
    "donate":      f"{BASE}/Screenshot 2026-03-25 at 6.48.19 AM.png",
    "blog":        f"{BASE}/Screenshot 2026-03-25 at 6.49.43 AM.png",
    "testimonials":f"{BASE}/Screenshot 2026-03-25 at 6.49.54 AM.png",
    "faq":         f"{BASE}/Screenshot 2026-03-25 at 6.50.06 AM.png",
}

# ── Presentation setup ───────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

W = prs.slide_width
H = prs.slide_height

# ── Helper functions ─────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_image(slide, img_path, left, top, width, height):
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, left, top, width, height)
    else:
        # placeholder box
        add_rect(slide, left, top, width, height, fill_color=RGBColor(0xCC,0xCC,0xCC))
        add_text(slide, f"[Image: {os.path.basename(img_path)}]",
                 left, top + height//2 - Pt(10), width, Pt(20),
                 font_size=10, align=PP_ALIGN.CENTER, color=GRAY)


def green_header_slide(title, subtitle=None):
    """Returns a slide with a full-width green header bar."""
    slide = prs.slides.add_slide(BLANK)
    # green top bar
    add_rect(slide, 0, 0, W, Inches(1.4), fill_color=GREEN)
    add_text(slide, title,
             Inches(0.5), Inches(0.15), W - Inches(1), Inches(0.9),
             font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.95), W - Inches(1), Inches(0.5),
                 font_size=16, color=RGBColor(0xCC, 0xFF, 0xCC), align=PP_ALIGN.LEFT)
    return slide


def bullet_points(slide, items, left, top, width, height,
                  font_size=18, color=DARK, bullet="•  ", line_spacing=1.3):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)

# full dark green background
add_rect(slide, 0, 0, W, H, fill_color=GREEN)

# white accent bar (bottom)
add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill_color=GREEN_LIGHT)

# Title
add_text(slide,
         "Rebuilding a Nonprofit's Digital Presence",
         Inches(0.7), Inches(1.6), W - Inches(1.4), Inches(1.0),
         font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide,
         "Using Conversational AI",
         Inches(0.7), Inches(2.55), W - Inches(1.4), Inches(0.8),
         font_size=38, bold=True, color=RGBColor(0xA8, 0xE6, 0xBF), align=PP_ALIGN.CENTER)

# Divider line
add_rect(slide, Inches(4.5), Inches(3.45), Inches(4.3), Inches(0.04),
         fill_color=GREEN_LIGHT)

# Subtitle
add_text(slide,
         "A Case Study of Asociación Deportiva Curundú",
         Inches(0.7), Inches(3.6), W - Inches(1.4), Inches(0.5),
         font_size=17, italic=True, color=RGBColor(0xCC, 0xFF, 0xCC), align=PP_ALIGN.CENTER)

# Author info
add_text(slide,
         "Connor Solvason  ·  FECON 390  ·  Prof. Daniel Egger  ·  April 2026",
         Inches(0.7), Inches(6.6), W - Inches(1.4), Inches(0.5),
         font_size=13, color=RGBColor(0x99, 0xCC, 0xAA), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Who is ADC & The Problem
# ════════════════════════════════════════════════════════════════════════════
slide = green_header_slide("Who is ADC — and What Did They Need?")

# Left column: org facts
add_text(slide, "Asociación Deportiva Curundú",
         Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.5),
         font_size=20, bold=True, color=GREEN)

bullet_points(slide, [
    "Founded 2014 in Curundú, Panama",
    "Nonprofit: youth soccer + education + nutrition",
    "Serves 70+ athletes and 60+ families",
    "Run by 3 founders — no full-time staff",
    "Active on Instagram (@adcurundu)",
], Inches(0.5), Inches(2.15), Inches(5.5), Inches(3.5), font_size=17)

# Right column: the problem
add_rect(slide, Inches(6.8), Inches(1.55), Inches(6.0), Inches(5.5),
         fill_color=LIGHT_BG)

add_text(slide, "The Problem",
         Inches(7.0), Inches(1.75), Inches(5.6), Inches(0.5),
         font_size=20, bold=True, color=GREEN)

bullet_points(slide, [
    "Outdated Squarespace site — looked amateur",
    "No budget for a web developer",
    "Applying for grants with no credible web presence",
    "Founders had no coding experience",
    "Goal: professional site that earns donor & funder trust",
], Inches(7.0), Inches(2.3), Inches(5.6), Inches(3.8), font_size=17, color=DARK)

# big $ callout
add_text(slide, "$0", Inches(0.5), Inches(5.8), Inches(2.0), Inches(0.8),
         font_size=40, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(slide, "development budget",
         Inches(0.5), Inches(6.3), Inches(2.5), Inches(0.5),
         font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Method
# ════════════════════════════════════════════════════════════════════════════
slide = green_header_slide("The Method: Conversational AI Development")

# Three columns: what / how / result
cols = [
    (Inches(0.4),  "What",  "Describe every feature in plain English — no code written by hand"),
    (Inches(4.6),  "How",   "Claude AI generated all HTML, CSS, JavaScript, and Python in response to natural language prompts"),
    (Inches(8.8),  "Result","A production-ready 8-page website deployed on GitHub Pages — built by a non-developer"),
]

for x, heading, body in cols:
    add_rect(slide, x, Inches(1.6), Inches(3.9), Inches(5.3), fill_color=LIGHT_BG)
    add_text(slide, heading,
             x + Inches(0.15), Inches(1.8), Inches(3.6), Inches(0.55),
             font_size=22, bold=True, color=GREEN)
    add_rect(slide, x + Inches(0.15), Inches(2.4), Inches(0.5), Inches(0.05),
             fill_color=GREEN_LIGHT)
    add_text(slide, body,
             x + Inches(0.15), Inches(2.6), Inches(3.6), Inches(3.8),
             font_size=16, color=DARK)

# Stats bar at bottom
add_rect(slide, 0, Inches(6.9), W, Inches(0.6), fill_color=GREEN)
stats = [
    ("48", "Prompts documented"),
    ("8",  "Pages built"),
    ("5",  "Development phases"),
    ("6",  "Weeks start to finish"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(0.5) + i * Inches(3.2)
    add_text(slide, num,   x, Inches(6.88), Inches(1.2), Inches(0.38),
             font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + Inches(1.1), Inches(6.98), Inches(2.0), Inches(0.35),
             font_size=11, color=RGBColor(0xCC, 0xFF, 0xCC), align=PP_ALIGN.LEFT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Before vs. After: Homepage Hero
# ════════════════════════════════════════════════════════════════════════════
slide = green_header_slide(
    "Before vs. After: Homepage",
    "One prompt → professional hero section with video background and unified navigation"
)

# The screenshot is already side-by-side (new left, old right)
add_image(slide, IMGS["hero"],
          Inches(0.4), Inches(1.55), W - Inches(0.8), Inches(5.1))

# Labels
add_rect(slide, Inches(0.4), Inches(6.65), Inches(6.1), Inches(0.4), fill_color=GREEN_LIGHT)
add_text(slide, "NEW — classroom video background, clean 5-item nav, ADC branding",
         Inches(0.5), Inches(6.67), Inches(6.0), Inches(0.35),
         font_size=11, bold=True, color=WHITE)

add_rect(slide, Inches(6.85), Inches(6.65), Inches(6.1), Inches(0.4),
         fill_color=RGBColor(0x99, 0x99, 0x99))
add_text(slide, "OLD — Squarespace default layout, cluttered navigation",
         Inches(6.95), Inches(6.67), Inches(6.0), Inches(0.35),
         font_size=11, bold=True, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Before vs. After: About Page
# ════════════════════════════════════════════════════════════════════════════
slide = green_header_slide(
    "Before vs. After: About Page",
    "Founder bios, impact statistics, and mission — all added in a few prompts"
)

add_image(slide, IMGS["stats"],
          Inches(0.4), Inches(1.55), W - Inches(0.8), Inches(5.1))

add_rect(slide, Inches(0.4), Inches(6.65), Inches(6.1), Inches(0.4), fill_color=GREEN_LIGHT)
add_text(slide, "NEW — animated counters: 70+ athletes, 60+ families, 5+ years, 11 wins",
         Inches(0.5), Inches(6.67), Inches(6.0), Inches(0.35),
         font_size=11, bold=True, color=WHITE)

add_rect(slide, Inches(6.85), Inches(6.65), Inches(6.1), Inches(0.4),
         fill_color=RGBColor(0x99, 0x99, 0x99))
add_text(slide, "OLD — basic Squarespace text block, no structure or data",
         Inches(6.95), Inches(6.67), Inches(6.0), Inches(0.35),
         font_size=11, bold=True, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Before vs. After: Participation Page
# ════════════════════════════════════════════════════════════════════════════
slide = green_header_slide(
    "Before vs. After: Participation Page",
    "Expanded from one registration link to six audience-specific pathways"
)

add_image(slide, IMGS["participate"],
          Inches(0.4), Inches(1.55), W - Inches(0.8), Inches(5.1))

add_rect(slide, Inches(0.4), Inches(6.65), Inches(6.1), Inches(0.4), fill_color=GREEN_LIGHT)
add_text(slide, "NEW — cards for athletes, volunteers, MET students, equipment donors, sponsors, partners",
         Inches(0.5), Inches(6.67), Inches(6.0), Inches(0.35),
         font_size=11, bold=True, color=WHITE)

add_rect(slide, Inches(6.85), Inches(6.65), Inches(6.1), Inches(0.4),
         fill_color=RGBColor(0x99, 0x99, 0x99))
add_text(slide, "OLD — single registration form for athletes only",
         Inches(6.95), Inches(6.67), Inches(6.0), Inches(0.35),
         font_size=11, bold=True, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Key Prompts
# ════════════════════════════════════════════════════════════════════════════
slide = green_header_slide("What a Prompt Actually Looks Like")

# Prompt 1
add_rect(slide, Inches(0.4), Inches(1.6), W - Inches(0.8), Inches(2.1),
         fill_color=RGBColor(0xE8, 0xF5, 0xE9))
add_rect(slide, Inches(0.4), Inches(1.6), Inches(0.12), Inches(2.1), fill_color=GREEN_LIGHT)
add_text(slide, "PROMPT (Phase 1 — Homepage)",
         Inches(0.65), Inches(1.65), W - Inches(1.2), Inches(0.35),
         font_size=11, bold=True, color=GREEN)
add_text(slide,
         '"Build a complete multi-page website for ADC, a nonprofit soccer and education organization in '
         'Curundú, Panama. Include a hero section with a video background, a clean navigation menu '
         'with these pages: Acerca de ADC, Calendario, Participar, Contáctanos, Blog, and Donar. '
         'Match ADC\'s green and black brand colors."',
         Inches(0.65), Inches(2.05), W - Inches(1.2), Inches(1.45),
         font_size=14, italic=True, color=DARK)

# Arrow
add_text(slide, "↓  Claude returned 400+ lines of production-ready HTML/CSS/JS",
         Inches(0.4), Inches(3.75), W - Inches(0.8), Inches(0.4),
         font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Prompt 2
add_rect(slide, Inches(0.4), Inches(4.2), W - Inches(0.8), Inches(1.85),
         fill_color=RGBColor(0xE8, 0xF5, 0xE9))
add_rect(slide, Inches(0.4), Inches(4.2), Inches(0.12), Inches(1.85), fill_color=GREEN_LIGHT)
add_text(slide, "PROMPT (Phase 4 — Grant Writer)",
         Inches(0.65), Inches(4.25), W - Inches(1.2), Inches(0.35),
         font_size=11, bold=True, color=GREEN)
add_text(slide,
         '"Build a private admin dashboard for ADC\'s founders. Include a password gate, a grant '
         'proposal writer powered by Claude AI, a to-do section, and Google Analytics links. '
         'Host it separately from the public site."',
         Inches(0.65), Inches(4.65), W - Inches(1.2), Inches(1.2),
         font_size=14, italic=True, color=DARK)

add_text(slide, "↓  Claude built a full backend: Cloudflare Worker + serverless API proxy + dashboard UI",
         Inches(0.4), Inches(6.1), W - Inches(0.8), Inches(0.4),
         font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — New Pages Built (Calendar + Donation)
# ════════════════════════════════════════════════════════════════════════════
slide = green_header_slide(
    "Pages Built From Scratch — That Didn't Exist Before",
    "Calendar, Donation, Blog, FAQ — each built in a single prompt session"
)

# Two images side by side
img_w = (W - Inches(1.1)) / 2
add_image(slide, IMGS["calendar"],
          Inches(0.4), Inches(1.6), img_w, Inches(4.8))
add_image(slide, IMGS["donate"],
          Inches(0.5) + img_w, Inches(1.6), img_w, Inches(4.8))

add_text(slide, "Styled event list with color-coded badges",
         Inches(0.4), Inches(6.45), img_w, Inches(0.4),
         font_size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "Tiered giving levels + fund transparency progress bars",
         Inches(0.5) + img_w, Inches(6.45), img_w, Inches(0.4),
         font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — AI-Powered Grant Writer
# ════════════════════════════════════════════════════════════════════════════
slide = green_header_slide(
    "The AI Grant Writer: Beyond the Public Website",
    "Phase 4 produced a private admin backend — the most technically complex feature"
)

# Left: explanation
add_text(slide, "What It Does",
         Inches(0.5), Inches(1.65), Inches(5.8), Inches(0.45),
         font_size=20, bold=True, color=GREEN)

bullet_points(slide, [
    "Private dashboard — password protected, separate from public site",
    "Founder types in a grant opportunity → Claude drafts a full proposal",
    "Auto-saves drafts to browser storage",
    "Tracks pending tasks, links to Google Analytics",
    "Cloudflare Worker routes API calls — API key never exposed",
], Inches(0.5), Inches(2.2), Inches(5.8), Inches(3.8), font_size=16)

# Right: why it matters
add_rect(slide, Inches(7.0), Inches(1.6), Inches(5.9), Inches(5.4),
         fill_color=GREEN)
add_text(slide, "Why This Matters",
         Inches(7.2), Inches(1.8), Inches(5.5), Inches(0.45),
         font_size=20, bold=True, color=WHITE)

bullet_points(slide, [
    "Grant writers charge $50–150/hr",
    "Small nonprofits often can't afford one",
    "ADC can now generate polished proposals instantly",
    "Same AI that built the site now runs inside it",
    "Live at: adc-grant-proxy.adcurundu.workers.dev",
], Inches(7.2), Inches(2.35), Inches(5.5), Inches(4.0),
   font_size=16, color=WHITE, bullet="→  ")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — By the Numbers
# ════════════════════════════════════════════════════════════════════════════
slide = green_header_slide("By the Numbers")

stats_big = [
    ("48",   "natural language\nprompts"),
    ("8",    "pages built\nfrom scratch"),
    ("$0",   "developer\ncost"),
    ("400+", "lines of code\nper prompt"),
]

box_w = (W - Inches(1.0)) / 4
for i, (num, label) in enumerate(stats_big):
    x = Inches(0.5) + i * box_w
    add_rect(slide, x, Inches(1.7), box_w - Inches(0.15), Inches(3.0),
             fill_color=LIGHT_BG if i % 2 == 0 else RGBColor(0xD7,0xEE,0xDD))
    add_text(slide, num,
             x, Inches(2.1), box_w - Inches(0.15), Inches(1.2),
             font_size=52, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(slide, label,
             x, Inches(3.3), box_w - Inches(0.15), Inches(1.1),
             font_size=15, color=GRAY, align=PP_ALIGN.CENTER)

# Phase breakdown bar
add_text(slide, "5 phases  →  Architecture  ·  Content  ·  Pages  ·  AI Features  ·  Documentation",
         Inches(0.5), Inches(5.0), W - Inches(1.0), Inches(0.5),
         font_size=15, color=DARK, align=PP_ALIGN.CENTER)

# Bottom quote
add_rect(slide, Inches(1.5), Inches(5.7), W - Inches(3.0), Inches(1.4),
         fill_color=GREEN)
add_text(slide,
         '"Every line of HTML, CSS, JavaScript, and Python was generated by Claude.\n'
         'Not a single character of code was written by hand."',
         Inches(1.7), Inches(5.85), W - Inches(3.4), Inches(1.1),
         font_size=15, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — What This Proves
# ════════════════════════════════════════════════════════════════════════════
slide = green_header_slide(
    "What This Proves About AI Today",
    "The bottleneck is no longer technical skill — it's knowing what you want"
)

points = [
    ("Democratizes professional development",
     "A non-coder produced a site that would cost $10,000–$20,000 from an agency — in 6 weeks, for free."),
    ("AI closes the resource gap for small orgs",
     "70% of nonprofits lack a dedicated web presence. Conversational AI removes the only barrier: cost."),
    ("The model is the developer",
     "Not a copilot, not autocomplete — Claude reasoned about design, UX, branding, and security end-to-end."),
    ("Economic disruption is already here",
     "Web dev, grant writing, content strategy — entire service categories are now accessible via natural language."),
]

for i, (heading, body) in enumerate(points):
    row = i // 2
    col = i % 2
    x = Inches(0.4) + col * Inches(6.45)
    y = Inches(1.65) + row * Inches(2.55)
    add_rect(slide, x, y, Inches(6.2), Inches(2.3),
             fill_color=LIGHT_BG if (i % 2 == 0) else RGBColor(0xD7,0xEE,0xDD))
    add_rect(slide, x, y, Inches(0.1), Inches(2.3), fill_color=GREEN_LIGHT)
    add_text(slide, heading,
             x + Inches(0.2), y + Inches(0.15), Inches(5.85), Inches(0.45),
             font_size=16, bold=True, color=GREEN)
    add_text(slide, body,
             x + Inches(0.2), y + Inches(0.65), Inches(5.85), Inches(1.5),
             font_size=14, color=DARK)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Q&A / Conclusion
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, W, H, fill_color=GREEN)
add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), fill_color=GREEN_LIGHT)

add_text(slide, "Thank You",
         Inches(0.7), Inches(1.4), W - Inches(1.4), Inches(0.9),
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.0), Inches(2.4), Inches(5.3), Inches(0.05),
         fill_color=GREEN_LIGHT)

add_text(slide, "Questions?",
         Inches(0.7), Inches(2.6), W - Inches(1.4), Inches(0.7),
         font_size=28, italic=True, color=RGBColor(0xCC, 0xFF, 0xCC), align=PP_ALIGN.CENTER)

links = [
    "Live site:   adcurundu.org",
    "GitHub:      csolv.github.io/my-website",
    "Instagram:   @adcurundu",
]
for i, link in enumerate(links):
    add_text(slide, link,
             Inches(3.5), Inches(3.6) + i * Inches(0.65), Inches(6.3), Inches(0.5),
             font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide,
         "Connor Solvason  ·  FECON 390  ·  Prof. Daniel Egger  ·  April 2026",
         Inches(0.7), Inches(6.7), W - Inches(1.4), Inches(0.4),
         font_size=12, color=RGBColor(0x99, 0xCC, 0xAA), align=PP_ALIGN.CENTER)

# ── Save ─────────────────────────────────────────────────────────────────────
out_path = f"{BASE}/ADC_AI_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
